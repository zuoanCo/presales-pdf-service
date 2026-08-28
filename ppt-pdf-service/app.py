"""PPT 模板填充 + 高保真 PDF 生成服务。

启动：uvicorn app:app --host 0.0.0.0 --port 8000
接口文档（Swagger UI）：http://localhost:8000/docs

PDF 转换使用本机 PowerPoint（COM 自动化）或 LibreOffice，排版样式与 PPT 完全一致。
"""
from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import uuid
from pathlib import Path

from fastapi import FastAPI, File, Form, HTTPException, UploadFile
from fastapi.responses import FileResponse
from pydantic import BaseModel, Field

from pptx_fill import extract_placeholders, load_and_fill
from pptx import Presentation

BASE_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = BASE_DIR / "templates"
OUTPUT_DIR = BASE_DIR / "output"
TEMPLATES_DIR.mkdir(exist_ok=True)
OUTPUT_DIR.mkdir(exist_ok=True)

app = FastAPI(
    title="PPT 模板 PDF 生成服务",
    description=(
        "把参数填入 PPT 模板中的 `{{占位符}}`，导出与 PPT 排版完全一致的 PDF。\n\n"
        "**使用流程**：① `GET /templates` 查看可用模板 → "
        "② `GET /placeholders/{模板名}` 查看模板需要的参数 → "
        "③ `POST /render` 传参数生成 PDF。\n\n"
        "**转换后端**：Windows + Office 环境自动使用 PowerPoint（保真度 100%）；"
        "其他环境使用 LibreOffice（需安装模板所用字体）。"
        "可用环境变量 `PDF_CONVERTER=powerpoint|libreoffice` 强制指定。"
    ),
    version="1.0.0",
)


class RenderRequest(BaseModel):
    """按模板名渲染请求。"""

    template: str = Field(
        default="售前模板.pptx",
        description="模板文件名，必须是 templates/ 目录下已有的 .pptx，可通过 GET /templates 查看",
        examples=["售前模板.pptx"],
    )
    params: dict[str, str] = Field(
        default_factory=dict,
        description=(
            "参数键值对：key 为占位符名（不含 {{ }}），value 为填入的内容。"
            "value 中可用 \\n 换行（自动拆成多个段落）。"
            "未提供的占位符会在 PDF 中保留 {{参数名}} 原样"
        ),
    )
    filename: str | None = Field(
        default=None,
        description="输出 PDF 的文件名（不含 .pdf 后缀），不填则自动生成",
        examples=["售前模板-递进测试"],
    )

    model_config = {
        "json_schema_extra": {
            "examples": [
                {
                    "template": "售前模板.pptx",
                    "filename": "售前模板-递进测试",
                    "params": {
                        "核心诊断": "1",
                        "战法": "2",
                        "市场阶段": "3",
                        "用户心智状态": "4",
                        "企业资源要求": "5",
                        "核心战略目标": "6",
                        "行业趋势": "7",
                        "用户原点": "8",
                        "品类切口": "9",
                        "战略打法": "10",
                        "价值载体": "11",
                        "原点用户": "12",
                        "价值方案": "13",
                        "视觉锤语言钉": "14",
                        "价值心智": "15",
                        "产品名称": "16",
                    },
                }
            ]
        }
    }


def _convert_to_pdf(pptx_path: Path, pdf_path: Path) -> None:
    """在子进程中调用 PowerPoint/LibreOffice 转换，失败时抛出带 stderr 的异常。"""
    result = subprocess.run(
        [sys.executable, str(BASE_DIR / "pdf_convert.py"), str(pptx_path), str(pdf_path)],
        capture_output=True,
        text=True,
        timeout=300,
    )
    if result.returncode != 0 or not pdf_path.exists():
        raise RuntimeError(f"PDF 转换失败: {result.stderr or result.stdout}")


def _render(template_path: Path, params: dict[str, str], out_name: str) -> Path:
    job = uuid.uuid4().hex[:8]
    with tempfile.TemporaryDirectory() as tmp:
        filled_pptx = Path(tmp) / f"filled_{job}.pptx"
        missing = load_and_fill(str(template_path), params, str(filled_pptx))
        if missing:
            # 缺失参数不阻断生成，但记录到日志，方便排查
            print(f"[warn] 模板 {template_path.name} 缺少参数: {missing}")
        pdf_path = OUTPUT_DIR / f"{out_name}.pdf"
        _convert_to_pdf(filled_pptx, pdf_path)
    return pdf_path


@app.get(
    "/templates",
    summary="列出可用模板",
    description="返回 templates/ 目录下所有 .pptx 模板文件名。新模板放入该目录即可使用，无需重启服务。",
)
def list_templates() -> list[str]:
    return sorted(p.name for p in TEMPLATES_DIR.glob("*.pptx"))


@app.get(
    "/placeholders/{template_name}",
    summary="查看模板需要哪些参数",
    description="扫描模板全文（含表格、组合形状），返回所有 `{{占位符}}` 的名字列表，即调用 /render 时 params 里可以填的 key。",
)
def get_placeholders(template_name: str):
    template_path = TEMPLATES_DIR / template_name
    if not template_path.exists():
        raise HTTPException(404, f"模板不存在: {template_name}")
    prs = Presentation(str(template_path))
    return {"template": template_name, "placeholders": extract_placeholders(prs)}


@app.post(
    "/render",
    summary="按模板名生成 PDF",
    description=(
        "用 templates/ 目录下已有的模板生成 PDF。\n\n"
        "- 31 页模板渲染约需 15 秒 ~ 2 分钟（PowerPoint 启动开销），请耐心等待；\n"
        "- 响应直接返回 PDF 文件，在 Swagger 页面点 **Download file** 下载；\n"
        "- 生成的 PDF 同时保留在服务器 output/ 目录下。"
    ),
)
def render(req: RenderRequest) -> FileResponse:
    template_path = TEMPLATES_DIR / req.template
    if not template_path.exists():
        raise HTTPException(404, f"模板不存在: {req.template}，可用 GET /templates 查看已有模板")
    out_name = req.filename or f"{template_path.stem}_{uuid.uuid4().hex[:6]}"
    try:
        pdf_path = _render(template_path, req.params, out_name)
    except Exception as e:
        raise HTTPException(500, str(e))
    return FileResponse(str(pdf_path), media_type="application/pdf", filename=f"{out_name}.pdf")


@app.post(
    "/render-upload",
    summary="上传模板文件生成 PDF",
    description=(
        "不上模板到服务器、直接上传一个 .pptx 文件进行渲染，适合临时测试新模板。\n\n"
        "- `file`：模板文件，内容用 `{{参数名}}` 标记可替换位置；\n"
        "- `params`：JSON 字符串，如 `{\"姓名\": \"张三\"}`。"
    ),
)
def render_upload(
    file: UploadFile = File(..., description="PPT 模板文件（.pptx）"),
    params: str = Form(default="{}", description="参数键值对的 JSON 字符串"),
) -> FileResponse:
    try:
        params_dict = json.loads(params)
    except json.JSONDecodeError:
        raise HTTPException(400, "params 必须是合法 JSON 字符串")

    with tempfile.TemporaryDirectory() as tmp:
        template_path = Path(tmp) / (file.filename or "template.pptx")
        template_path.write_bytes(file.file.read())
        out_name = f"{template_path.stem}_{uuid.uuid4().hex[:6]}"
        try:
            pdf_path = _render(template_path, params_dict, out_name)
        except Exception as e:
            raise HTTPException(500, str(e))
        # FileResponse 在后台读文件，TemporaryDirectory 退出即删，这里先读入内存路径副本
        data = pdf_path.read_bytes()
        final_path = OUTPUT_DIR / f"{out_name}.pdf"
        final_path.write_bytes(data)
    return FileResponse(str(final_path), media_type="application/pdf", filename=f"{out_name}.pdf")


if __name__ == "__main__":
    import uvicorn

    uvicorn.run(app, host="0.0.0.0", port=8000)
