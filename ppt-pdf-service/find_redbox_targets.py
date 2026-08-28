"""找出每个红框（FF0000 边框矩形）覆盖的文字内容。"""
from pptx import Presentation

prs = Presentation(r"C:\Users\15944\Desktop\售前模板.pptx")


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def is_red_box(shape):
    try:
        return shape.line.color.type == 1 and str(shape.line.color.rgb) == "FF0000"
    except Exception:
        return False


def box(shape):
    return (shape.left, shape.top, shape.left + shape.width, shape.top + shape.height)


def overlap(b1, b2):
    """返回 b1 与 b2 的相交面积占 b1 面积的比例"""
    x1, y1 = max(b1[0], b2[0]), max(b1[1], b2[1])
    x2, y2 = min(b1[2], b2[2]), min(b1[3], b2[3])
    if x2 <= x1 or y2 <= y1:
        return 0.0
    inter = (x2 - x1) * (y2 - y1)
    area1 = (b1[2] - b1[0]) * (b1[3] - b1[1])
    return inter / area1 if area1 else 0.0


for si, slide in enumerate(prs.slides, 1):
    shapes = list(iter_shapes(slide.shapes))
    reds = [s for s in shapes if is_red_box(s)]
    if not reds:
        continue
    print(f"\n===== slide {si}：{len(reds)} 个红框 =====")
    for red in reds:
        rb = box(red)
        print(f"\n  红框 {red.name!r} @ left={rb[0]}, top={rb[1]}, w={red.width}, h={red.height}")
        for s in shapes:
            if s is red or not s.has_text_frame or not s.text_frame.text.strip():
                continue
            ratio = overlap(box(s), rb)
            if ratio > 0.3:
                text = s.text_frame.text.replace("\n", " | ")[:120]
                print(f"    -> 覆盖 {ratio:.0%} [{s.name}] {text!r}")
