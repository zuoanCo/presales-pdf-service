"""检查模板：列出所有形状的边框颜色和文字，找出红框标记的参数位置。"""
import sys
from pptx import Presentation

path = sys.argv[1] if len(sys.argv) > 1 else r"C:\Users\15944\Desktop\售前模板.pptx"
prs = Presentation(path)
print(f"slides: {len(prs.slides)}  size: {prs.slide_width} x {prs.slide_height}")


def iter_shapes(shapes, prefix=""):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes, prefix + "  ")
        else:
            yield prefix, shape


for si, slide in enumerate(prs.slides, 1):
    print(f"\n===== slide {si} =====")
    for prefix, shape in iter_shapes(slide.shapes):
        color = None
        try:
            line = shape.line
            if line.color and line.color.type is not None:
                if line.color.type == 1:
                    color = str(line.color.rgb)
                else:
                    color = f"theme:{line.color.theme_color}"
        except Exception:
            pass
        text = ""
        if shape.has_text_frame:
            text = shape.text_frame.text.replace("\n", " | ")[:80]
        if color or text:
            print(f"{prefix}[{shape.shape_type}] name={shape.name!r} line={color} text={text!r}")
