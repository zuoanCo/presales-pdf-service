"""生成一个带占位符的示例模板，用于端到端测试。

演示两类典型版式：封面（标题+副标题占位符）和内容页（正文+表格占位符）。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE_DIR = Path(__file__).resolve().parent
TEMPLATES_DIR = BASE_DIR / "templates"
TEMPLATES_DIR.mkdir(exist_ok=True)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

# ---- 第 1 页：封面 ----
slide = prs.slides.add_slide(blank)

# 背景色块
bg = slide.shapes.add_shape(1, 0, 0, prs.slide_width, Inches(2.6))  # 1 = 矩形
bg.fill.solid()
bg.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)
bg.line.fill.background()

title = slide.shapes.add_textbox(Inches(1), Inches(0.9), Inches(11.3), Inches(1.2))
tf = title.text_frame
tf.text = "{{标题}}"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.runs[0]
r.font.size = Pt(44)
r.font.bold = True
r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
r.font.name = "微软雅黑"

sub = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(11.3), Inches(0.8))
tf = sub.text_frame
tf.text = "汇报人：{{姓名}}    部门：{{部门}}"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(24)
p.runs[0].font.name = "微软雅黑"

date_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.3), Inches(0.6))
tf = date_box.text_frame
tf.text = "{{日期}}"
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.runs[0].font.size = Pt(18)
p.runs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
p.runs[0].font.name = "微软雅黑"

# ---- 第 2 页：内容页（文本 + 表格） ----
slide2 = prs.slides.add_slide(blank)
head = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(0.9))
tf = head.text_frame
tf.text = "{{标题}} — 关键数据"
p = tf.paragraphs[0]
p.runs[0].font.size = Pt(32)
p.runs[0].font.bold = True
p.runs[0].font.color.rgb = RGBColor(0x1F, 0x4E, 0x79)
p.runs[0].font.name = "微软雅黑"

body = slide2.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(12), Inches(1.5))
tf = body.text_frame
tf.word_wrap = True
tf.text = "本季度由 {{姓名}} 负责的项目整体进度 {{进度}}，核心指标完成情况如下："
p = tf.paragraphs[0]
p.runs[0].font.size = Pt(18)
p.runs[0].font.name = "微软雅黑"

tbl_shape = slide2.shapes.add_table(3, 3, Inches(0.5), Inches(2.6), Inches(8), Inches(2.4))
tbl = tbl_shape.table
headers = ["指标", "目标", "实际"]
rows = [["销售额", "1000万", "{{销售额}}"], ["新客户", "50", "{{新客户数}}"]]
for j, h in enumerate(headers):
    cell = tbl.cell(0, j)
    cell.text = h
    cell.text_frame.paragraphs[0].runs[0].font.bold = True
    cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
    cell.text_frame.paragraphs[0].runs[0].font.name = "微软雅黑"
for i, row in enumerate(rows, start=1):
    for j, v in enumerate(row):
        cell = tbl.cell(i, j)
        cell.text = v
        cell.text_frame.paragraphs[0].runs[0].font.size = Pt(16)
        cell.text_frame.paragraphs[0].runs[0].font.name = "微软雅黑"

out = TEMPLATES_DIR / "汇报模板.pptx"
prs.save(str(out))
print(f"模板已生成: {out}")
