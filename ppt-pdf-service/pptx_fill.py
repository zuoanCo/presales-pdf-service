"""PPT 模板占位符填充模块。

约定：模板中用 {{参数名}} 作为占位符，例如 {{姓名}}、{{日期}}。
填充时完整保留 PPT 原有的字体、字号、颜色、加粗等格式。
"""
from __future__ import annotations

import copy
import re
from typing import Mapping

from pptx import Presentation
from pptx.oxml.ns import qn

PLACEHOLDER_RE = re.compile(r"\{\{\s*([^{}]+?)\s*\}\}")


def extract_placeholders(prs: Presentation) -> list[str]:
    """扫描整个演示文稿，返回出现过的占位符名（去重、保持出现顺序）。"""
    found: list[str] = []

    def _collect(text: str) -> None:
        for m in PLACEHOLDER_RE.finditer(text):
            name = m.group(1)
            if name not in found:
                found.append(name)

    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                _collect(shape.text_frame.text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _collect(cell.text_frame.text)
    return found


def fill_presentation(prs: Presentation, params: Mapping[str, str]) -> list[str]:
    """把 params 中的值填入所有占位符，返回未被提供的占位符名列表。"""
    used: set[str] = set()
    missing: list[str] = []

    def _replace_text(text: str) -> str:
        def _sub(m: re.Match) -> str:
            name = m.group(1)
            if name in params:
                used.add(name)
                return str(params[name])
            if name not in missing:
                missing.append(name)
            return m.group(0)  # 未提供的参数保留原样

        return PLACEHOLDER_RE.sub(_sub, text)

    for slide in prs.slides:
        for shape in _iter_shapes(slide.shapes):
            if shape.has_text_frame:
                _fill_text_frame(shape.text_frame, _replace_text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        _fill_text_frame(cell.text_frame, _replace_text)
    return missing


def _iter_shapes(shapes):
    """递归遍历形状，展开组合形状（group）。"""
    for shape in shapes:
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            yield from _iter_shapes(shape.shapes)
        else:
            yield shape


def _fill_text_frame(text_frame, replace_text) -> None:
    for paragraph in text_frame.paragraphs:
        _fill_paragraph(paragraph, replace_text)


def _fill_paragraph(paragraph, replace_text) -> None:
    """段落级替换：占位符可能被拆到多个 run 中，需要在段落整文本上替换。

    策略：如果替换只影响某一个 run，就只改那个 run 的文本（最大程度保留格式）；
    如果占位符横跨多个 run，则把替换后的整段文本写入第一个含占位符的 run，
    清空其余被波及的 run（格式继承第一个 run，这是不可避免的取舍）。
    """
    runs = paragraph.runs
    if not runs:
        return

    full = "".join(r.text for r in runs)
    if not PLACEHOLDER_RE.search(full):
        return

    # 先尝试：逐 run 替换（占位符完整落在单个 run 内的常见情况）
    changed = False
    for run in runs:
        if PLACEHOLDER_RE.search(run.text):
            run.text = replace_text(run.text)
            changed = True
    if changed:
        _split_multiline_runs(paragraph)
        return

    # 占位符横跨多个 run：定位区间，合并替换
    spans = []  # (start, end) 每个 run 在 full 中的位置
    pos = 0
    for r in runs:
        spans.append((pos, pos + len(r.text)))
        pos += len(r.text)

    new_full = replace_text(full)
    if new_full == full:
        return

    first_idx = None
    for i, (s, e) in enumerate(spans):
        seg = full[s:e]
        # 找第一个与占位符区间相交的 run
        for m in PLACEHOLDER_RE.finditer(full):
            if s < m.end() and m.start() < e:
                first_idx = i
                break
        if first_idx is not None:
            break
    if first_idx is None:
        first_idx = 0

    runs[first_idx].text = new_full
    for i, r in enumerate(runs):
        if i != first_idx:
            r.text = ""


def _split_multiline_runs(paragraph) -> None:
    """参数值含换行符时，把 run 拆成多个段落（克隆原段落/字符格式）。

    PowerPoint 里换行是新段落，直接在 run 里塞 \\n 不会被正确渲染。
    """
    for run in paragraph.runs:
        if "\n" not in run.text:
            continue
        lines = run.text.split("\n")
        run.text = lines[0]
        p_el = paragraph._p
        anchor = p_el
        for line in lines[1:]:
            new_p = copy.deepcopy(p_el)
            # 只保留第一个 run，去掉其余 run 和手动换行
            rs = new_p.findall(qn("a:r"))
            for extra in rs[1:]:
                new_p.remove(extra)
            for br in new_p.findall(qn("a:br")):
                new_p.remove(br)
            rs[0].find(qn("a:t")).text = line
            anchor.addnext(new_p)
            anchor = new_p
        # 原段落的 pPr/剩余 run 不动，后续段落已克隆
        return


def load_and_fill(template_path: str, params: Mapping[str, str], out_path: str) -> list[str]:
    """加载模板、填充参数、另存为新 pptx，返回缺失参数列表。"""
    prs = Presentation(template_path)
    missing = fill_presentation(prs, params)
    prs.save(out_path)
    return missing
