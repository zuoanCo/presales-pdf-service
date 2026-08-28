# -*- coding: utf-8 -*-
"""PHP 轻量版预处理（只需运行一次）：

1. 从 templates/售前模板.pptx 提取每个参数框的几何位置与字体样式 -> php/assets/config.json
2. 生成「参数区留空」的模板副本 -> php/assets/blank.pptx
3. 后续用 WPS 把 blank.pptx 转 PDF，再栅格化为每页 PNG 背景图
"""
import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

BASE = Path(__file__).resolve().parent
TEMPLATE = BASE / "templates" / "售前模板.pptx"
OUT = BASE / "php" / "assets"
OUT.mkdir(parents=True, exist_ok=True)

# (页码, 形状名, 参数名, 后缀)  后缀用于标题这类「参数+固定文字」的框
TARGETS = [
    (27, "行1字0", "战法", ""),
    (27, "行1字1", "市场阶段", ""),
    (27, "行1字2", "用户心智状态", ""),
    (27, "行1字3", "企业资源要求", ""),
    (27, "行1字4", "核心战略目标", ""),
    (27, "文本框 3", "核心诊断", ""),
    (28, "内容-① 行业趋势", "行业趋势", ""),
    (28, "内容-② 用户原点", "用户原点", ""),
    (29, "内容-③ 品类切口", "品类切口", ""),
    (29, "内容-④ 战略打法", "战略打法", ""),
    (30, "文本框 53", "产品名称", "三品合一判断"),
    (30, "文本框 46", "原点用户", ""),
    (30, "TextBox 16", "价值载体", ""),
    (30, "文本框 52", "价值心智", ""),
    (30, "文本框 50", "视觉锤语言钉", ""),
    (30, "文本框 48", "价值方案", ""),
]

EMU_PER_PT = 12700


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def body_pr_info(tf):
    bodyPr = tf._txBody.find(qn("a:bodyPr"))
    def i(attr, default):
        v = bodyPr.get(attr)
        return int(v) if v is not None else default
    return {
        "lIns": i("lIns", 91440), "tIns": i("tIns", 45720),
        "rIns": i("rIns", 91440), "bIns": i("bIns", 45720),
        "anchor": bodyPr.get("anchor", "t"),
        "wrap": bodyPr.get("wrap", "square"),
    }


def para_info(p):
    """段落的行距（倍数，None 表示默认）与对齐"""
    pPr = p._pPr
    spacing = None
    align = None
    if pPr is not None:
        lnSpc = pPr.find(qn("a:lnSpc"))
        if lnSpc is not None:
            pct = lnSpc.find(qn("a:spcPct"))
            pts = lnSpc.find(qn("a:spcPts"))
            if pct is not None:
                spacing = int(pct.get("val")) / 100000.0
            elif pts is not None:
                spacing = ("pt", int(pts.get("val")) / 100.0)
        align = pPr.get("algn")
    return spacing, align


def run_style(run):
    f = run.font
    color = None
    try:
        if f.color and f.color.type == 1:
            color = str(f.color.rgb)
    except Exception:
        pass
    # 东亚字体要读 <a:ea>，python-pptx 的 font.name 只读 latin
    rPr = run._r.find(qn("a:rPr"))
    ea = None
    if rPr is not None:
        ea_el = rPr.find(qn("a:ea"))
        if ea_el is not None:
            ea = ea_el.get("typeface")
    return {
        "font": ea or f.name,
        "size_pt": f.size.pt if f.size else None,
        "bold": f.bold,
        "color": color,
    }


prs = Presentation(str(TEMPLATE))
slide_w, slide_h = prs.slide_width, prs.slide_height
params_out = []

for slide_no, shape_name, param, suffix in TARGETS:
    slide = prs.slides[slide_no - 1]
    shape = next(s for s in iter_shapes(slide.shapes) if s.name == shape_name)
    tf = shape.text_frame
    # 用字符数最多的 run 作为格式原型（与 prepare_template 同一策略）
    runs = [r for p in tf.paragraphs for r in p.runs]
    proto = max(runs, key=lambda r: len(r.text)) if runs else None
    style = run_style(proto) if proto else {}
    spacing, align = para_info(tf.paragraphs[0])
    bp = body_pr_info(tf)
    params_out.append({
        "name": param,
        "suffix": suffix,
        "page": slide_no,
        "x_pt": round(shape.left / EMU_PER_PT, 2),
        "y_pt": round(shape.top / EMU_PER_PT, 2),
        "w_pt": round(shape.width / EMU_PER_PT, 2),
        "h_pt": round(shape.height / EMU_PER_PT, 2),
        "insets_pt": {
            "l": round(bp["lIns"] / EMU_PER_PT, 2), "t": round(bp["tIns"] / EMU_PER_PT, 2),
            "r": round(bp["rIns"] / EMU_PER_PT, 2), "b": round(bp["bIns"] / EMU_PER_PT, 2),
        },
        "anchor": bp["anchor"],
        "align": align or "l",
        "line_spacing": spacing,
        **style,
    })

    # 留空参数区：整框清空（标题框只清产品名部分由 suffix 处理——整框清空，后缀由 PHP 补画）
    for p in tf.paragraphs[1:]:
        p._p.getparent().remove(p._p)
    p0 = tf.paragraphs[0]
    for r in list(p0.runs):
        r._r.getparent().remove(r._r)
    # 保留一个空 run 占位，避免 bodyPr 结构异常
    p0._p.append(copy.deepcopy(p0._p)) if False else None
    from pptx.oxml.ns import qn as _qn
    r_el = p0._p.makeelement(_qn("a:r"), {})
    rPr_el = p0._p.makeelement(_qn("a:rPr"), {})
    t_el = p0._p.makeelement(_qn("a:t"), {})
    t_el.text = ""
    r_el.append(rPr_el)
    r_el.append(t_el)
    p0._p.append(r_el)

config = {
    "page_w_pt": round(slide_w / EMU_PER_PT, 2),
    "page_h_pt": round(slide_h / EMU_PER_PT, 2),
    "pages": len(prs.slides),
    "params": params_out,
}
(OUT / "config.json").write_text(json.dumps(config, ensure_ascii=False, indent=2), encoding="utf-8")
prs.save(str(OUT / "blank.pptx"))

print(f"页面: {config['page_w_pt']}x{config['page_h_pt']} pt, 共 {config['pages']} 页")
for p in params_out:
    print(f"  P{p['page']:>2} {p['name']:<8} font={p.get('font')} size={p.get('size_pt')} bold={p.get('bold')} color={p.get('color')} spacing={p['line_spacing']} align={p['align']}")
print(f"\n-> {OUT / 'config.json'}")
print(f"-> {OUT / 'blank.pptx'}")
