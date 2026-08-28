"""把 售前模板.pptx 红框标记的内容转成 {{占位符}}，并删除红框标记本身。

输出到 templates/售前模板.pptx，原文件不改动。
"""
from pathlib import Path

from pptx import Presentation

SRC = r"C:\Users\15944\Desktop\售前模板.pptx"
OUT = Path(__file__).resolve().parent / "templates" / "售前模板.pptx"

# (幻灯片页码从1开始, 形状名, 占位符名)
# 均为红框覆盖的内容；整体替换为占位符，格式保留第一个 run 的格式
TARGETS = [
    # slide 27：战法表格中被标记的一整行（5 个单元格）+ 核心诊断总结
    (27, "行1字0", "战法"),
    (27, "行1字1", "市场阶段"),
    (27, "行1字2", "用户心智状态"),
    (27, "行1字3", "企业资源要求"),
    (27, "行1字4", "核心战略目标"),
    (27, "文本框 3", "核心诊断"),
    # slide 28
    (28, "内容-① 行业趋势", "行业趋势"),
    (28, "内容-② 用户原点", "用户原点"),
    # slide 29
    (29, "内容-③ 品类切口", "品类切口"),
    (29, "内容-④ 战略打法", "战略打法"),
    # slide 30
    (30, "文本框 46", "原点用户"),
    (30, "TextBox 16", "价值载体"),
    (30, "文本框 52", "价值心智"),
    (30, "文本框 50", "视觉锤语言钉"),
    (30, "文本框 48", "价值方案"),
]

# slide 30 标题只替换产品名部分：'亲密陪伴机器人三品合一判断' -> '{{产品名称}}三品合一判断'
PARTIAL = [(30, "文本框 53", "亲密陪伴机器人", "{{产品名称}}")]


def iter_shapes(shapes):
    for shape in shapes:
        if shape.shape_type == 6:
            yield from iter_shapes(shape.shapes)
        else:
            yield shape


def set_placeholder(text_frame, placeholder: str) -> None:
    """整个文本框替换为占位符。

    格式原型选「字符数最多的 run」（正文 run），而不是第一个 run——
    这些文本框常以加粗标签开头（如「核心用户：」），用第一个 run 会把整段内容
    都变成加粗，视觉上比原设计重。正文 run 的格式更接近原设计的主视觉。
    """
    paras = text_frame.paragraphs
    all_runs = [r for p in paras for r in p.runs]
    if all_runs:
        proto = max(all_runs, key=lambda r: len(r.text))
        proto.text = placeholder
        proto_p = proto._r.getparent()
        # 删除同段其余 run
        for r in proto_p.findall(proto._r.tag):
            if r is not proto._r:
                proto_p.remove(r)
        # 删除其余所有段落
        for p in paras:
            if p._p is not proto_p:
                p._p.getparent().remove(p._p)
    else:
        paras[0].text = placeholder
        for p in paras[1:]:
            p._p.getparent().remove(p._p)


def is_red_box(shape) -> bool:
    """直接读 XML 判断边框颜色。

    不能用 shape.line.color——python-pptx 访问 shape.line 时会惰性创建空的
    <a:ln/> 元素，空边框按主题默认渲染成黑框，会把模板改坏（已踩坑）。
    """
    from pptx.oxml.ns import qn

    spPr = shape._element.find(qn("p:spPr"))
    if spPr is None:
        return False
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        return False
    clr = ln.find(f".//{qn('a:srgbClr')}")
    return clr is not None and clr.get("val", "").upper() == "FF0000"


prs = Presentation(SRC)

# 1. 替换占位符
replaced, missed = 0, []
for slide_no, shape_name, param in TARGETS:
    slide = prs.slides[slide_no - 1]
    target = next((s for s in iter_shapes(slide.shapes) if s.name == shape_name), None)
    if target is None:
        missed.append((slide_no, shape_name))
        continue
    set_placeholder(target.text_frame, "{{" + param + "}}")
    replaced += 1

for slide_no, shape_name, old, new in PARTIAL:
    slide = prs.slides[slide_no - 1]
    target = next((s for s in iter_shapes(slide.shapes) if s.name == shape_name), None)
    done = False
    if target is not None:
        for p in target.text_frame.paragraphs:
            for r in p.runs:
                if old in r.text:
                    r.text = r.text.replace(old, new)
                    done = True
    if done:
        replaced += 1
    else:
        missed.append((slide_no, shape_name + " (partial)"))

# 2. 删除所有红框标记
removed = 0
for slide in prs.slides:
    for shape in list(iter_shapes(slide.shapes)):
        if is_red_box(shape):
            shape._element.getparent().remove(shape._element)
            removed += 1

OUT.parent.mkdir(exist_ok=True)
prs.save(str(OUT))
print(f"替换占位符 {replaced} 处，删除红框 {removed} 个 -> {OUT}")
if missed:
    print("未找到的形状:", missed)
